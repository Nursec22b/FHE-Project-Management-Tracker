#!/usr/bin/env python3
"""Generate the FHE Project Board pitch presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# ============ COLORS ============
DARK_BLUE = RGBColor(0x02, 0x6A, 0xA7)
PRIMARY = RGBColor(0x00, 0x79, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF5, 0xF7)
TEXT_DARK = RGBColor(0x17, 0x2B, 0x4D)
TEXT_SEC = RGBColor(0x5E, 0x6C, 0x84)
GREEN = RGBColor(0x61, 0xBD, 0x4F)
RED = RGBColor(0xEB, 0x5A, 0x46)
YELLOW = RGBColor(0xF2, 0xD6, 0x00)
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT_TEAL = RGBColor(0x00, 0xAE, 0xCC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, w, h, size=18, color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet(tf, text, size=18, color=TEXT_DARK, bold=False, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.space_before = space_before
    return p

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Accent bar top
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

# Left accent stripe
add_shape(slide, Inches(0), Inches(0), Inches(0.25), Inches(7.5), PRIMARY)

# Decorative circles (top right)
c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-0.5), Inches(2), Inches(2))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x00, 0x56, 0x8A); c1.line.fill.background()
c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.3), Inches(0.3), Inches(2.5), Inches(2.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x01, 0x45, 0x70); c2.line.fill.background()

# Logo icon (kanban board shape)
add_rounded(slide, Inches(1.2), Inches(1.8), Inches(0.65), Inches(0.85), PRIMARY)
add_rounded(slide, Inches(1.95), Inches(1.8), Inches(0.65), Inches(0.6), ACCENT_TEAL)

# Title
add_text(slide, 'FHE Project Board', Inches(1.1), Inches(2.9), Inches(8), Inches(1),
         size=52, color=WHITE, bold=True)

# Subtitle
add_text(slide, 'Why Pay Rent When You Can Own?', Inches(1.1), Inches(3.8), Inches(8), Inches(0.6),
         size=26, color=ACCENT_TEAL, bold=False)

# Tagline
add_text(slide, 'An Internal Project Management Tool Built By Engineers, For Engineers',
         Inches(1.1), Inches(4.5), Inches(8), Inches(0.5), size=16, color=TEXT_SEC)

# Bottom bar
add_shape(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), RGBColor(0x01, 0x3A, 0x5E))
add_text(slide, 'Florida Horizon Engineering  |  February 2026  |  Confidential',
         Inches(1), Inches(7.05), Inches(11), Inches(0.4), size=12, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.06), RED)

add_text(slide, 'Houston, We Have a Subscription', Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
         size=38, color=WHITE, bold=True)

# Money flying away graphic - dollar signs
for i, x in enumerate([2.0, 5.5, 9.0]):
    y = 1.8
    box = add_rounded(slide, Inches(x), Inches(y), Inches(2.5), Inches(1.5), RGBColor(0xFC, 0xEA, 0xE8))
    add_text(slide, '$', Inches(x+0.15), Inches(y+0.1), Inches(0.5), Inches(1.3),
             size=48, color=RED, bold=True)

add_text(slide, '$62.50', Inches(2.5), Inches(2.0), Inches(1.5), Inches(0.5),
         size=28, color=RED, bold=True)
add_text(slide, '/month', Inches(2.5), Inches(2.5), Inches(1.5), Inches(0.3),
         size=14, color=TEXT_SEC)

add_text(slide, '$750', Inches(6.0), Inches(2.0), Inches(1.5), Inches(0.5),
         size=28, color=RED, bold=True)
add_text(slide, '/year', Inches(6.0), Inches(2.5), Inches(1.5), Inches(0.3),
         size=14, color=TEXT_SEC)

add_text(slide, '$3,750', Inches(9.4), Inches(2.0), Inches(1.5), Inches(0.5),
         size=28, color=RED, bold=True)
add_text(slide, '/ 5 years', Inches(9.4), Inches(2.5), Inches(1.5), Inches(0.3),
         size=14, color=TEXT_SEC)

# Labels
add_text(slide, '5 Trello Premium licenses x $12.50/user/month',
         Inches(1), Inches(3.6), Inches(11), Inches(0.4), size=18, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(1.5), Inches(4.2), Inches(10.333), Inches(0.02), RGBColor(0xDF, 0xE1, 0xE6))

# Scaling pain
box = add_rounded(slide, Inches(1), Inches(4.5), Inches(11.333), Inches(1.2), RGBColor(0xFF, 0xF4, 0xE5))
txBox = add_text(slide, '', Inches(1.4), Inches(4.55), Inches(10.5), Inches(1.1), size=16, color=TEXT_DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = "And it only gets worse as we grow..."
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0xD2, 0x90, 0x34)
add_bullet(tf, "10 users = $125/month = $1,500/year", size=16, color=TEXT_DARK)
add_bullet(tf, "20 users = $250/month = $3,000/year", size=16, color=TEXT_DARK)
add_bullet(tf, "20 users over 5 years = $15,000 paid to Atlassian", size=16, color=RED, bold=True)

# Humor
add_text(slide, '"Every time we add a team member, Atlassian sends us a thank-you card."',
         Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
         size=16, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: WHAT TRELLO CAN'T DO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.06), YELLOW)

add_text(slide, "What Trello CAN'T Do", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
         size=38, color=WHITE, bold=True)

# Pain points as cards
pain_points = [
    ("No Email Automation", "We need incoming emails to auto-create tasks. Trello says: 'upgrade to Enterprise.'", RED),
    ("Per-Seat Pricing", "Every new hire = more $$$. Growth shouldn't cost us more in software.", RGBColor(0xD2, 0x90, 0x34)),
    ("Limited Customization", "Our workflows don't fit Trello's mold. We adapt to them instead of them to us.", PRIMARY),
    ("Their Data, Their Rules", "Our project data lives on Atlassian's servers. We have zero control.", RGBColor(0x89, 0x60, 0x9E)),
]

for i, (title, desc, color) in enumerate(pain_points):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.7 + row * 2.4

    # Card
    card = add_rounded(slide, Inches(x), Inches(y), Inches(5.6), Inches(2.0), LIGHT_BG)

    # Color accent bar on left
    add_shape(slide, Inches(x), Inches(y), Inches(0.12), Inches(2.0), color)

    # X icon
    add_text(slide, '\u2717', Inches(x + 0.35), Inches(y + 0.2), Inches(0.5), Inches(0.5),
             size=28, color=RED, bold=True)

    add_text(slide, title, Inches(x + 0.85), Inches(y + 0.25), Inches(4.5), Inches(0.4),
             size=20, color=TEXT_DARK, bold=True)
    add_text(slide, desc, Inches(x + 0.85), Inches(y + 0.75), Inches(4.5), Inches(0.9),
             size=15, color=TEXT_SEC)

# Bottom quote
add_text(slide, '"We\'re renting a tool when we could own one."',
         Inches(1), Inches(6.5), Inches(11.333), Inches(0.5),
         size=20, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

add_text(slide, 'Introducing: FHE Project Board', Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
         size=40, color=WHITE, bold=True)
add_text(slide, 'Built BY us, FOR us. Zero subscription fees. Unlimited users. Forever.',
         Inches(0.8), Inches(1.2), Inches(10), Inches(0.4), size=18, color=ACCENT_TEAL)

features = [
    ("\u2611 Kanban Boards", "Drag & drop project management\njust like Trello", PRIMARY),
    ("\u2611 Email Automation", "Incoming emails automatically\ncreate task cards", GREEN),
    ("\u2611 Unlimited Users", "Add 5, 50, or 500 users.\nSame cost: $0 per seat", ACCENT_TEAL),
    ("\u2611 Custom Workflows", "Tailored to FHE's processes,\nnot Trello's templates", RGBColor(0xCD, 0x5A, 0x91)),
    ("\u2611 Data Ownership", "Our data stays on our servers.\nFull control, full privacy", RGBColor(0x89, 0x60, 0x9E)),
    ("\u2611 Mobile Ready", "Works on any device, anywhere.\nPWA with offline support", RGBColor(0xD2, 0x90, 0x34)),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.2
    y = 2.1 + row * 2.5

    card = add_rounded(slide, Inches(x), Inches(y), Inches(3.8), Inches(2.1), RGBColor(0x14, 0x28, 0x3C))
    add_shape(slide, Inches(x), Inches(y + 2.0), Inches(3.8), Inches(0.1), color)

    add_text(slide, title, Inches(x + 0.3), Inches(y + 0.25), Inches(3.2), Inches(0.4),
             size=20, color=WHITE, bold=True)
    add_text(slide, desc, Inches(x + 0.3), Inches(y + 0.8), Inches(3.2), Inches(1.0),
             size=15, color=TEXT_SEC)

# ============================================================
# SLIDE 5: LIVE DEMO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

# Big center circle with play icon
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(1.5), Inches(3), Inches(3))
circle.fill.solid()
circle.fill.fore_color.rgb = PRIMARY
circle.line.fill.background()

# Play triangle
add_text(slide, '\u25B6', Inches(5.7), Inches(1.8), Inches(2), Inches(2.5),
         size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Don't Take My Word For It", Inches(1), Inches(4.8), Inches(11.333), Inches(0.7),
         size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, 'LIVE DEMO TIME', Inches(1), Inches(5.5), Inches(11.333), Inches(0.5),
         size=24, color=ACCENT_TEAL, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, 'Yes, it actually works. No, it\'s not a PowerPoint animation.',
         Inches(1), Inches(6.3), Inches(11.333), Inches(0.4),
         size=16, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: THE MONEY SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.06), GREEN)

add_text(slide, 'Show Me The Money', Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
         size=38, color=WHITE, bold=True)

# LEFT: Trello (red)
add_rounded(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.2), RGBColor(0xFC, 0xEA, 0xE8))
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(0.6), RED)
add_text(slide, 'Trello Premium (5 years, 20 users)', Inches(1.0), Inches(1.65), Inches(5.2), Inches(0.5),
         size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, '$15,000', Inches(1.2), Inches(2.5), Inches(5), Inches(0.8),
         size=56, color=RED, bold=True, align=PP_ALIGN.CENTER)

txBox = add_text(slide, '', Inches(1.5), Inches(3.5), Inches(4.5), Inches(2.0), size=15, color=TEXT_DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = "$250/month at 20 users"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = TEXT_DARK
add_bullet(tf, "$3,000/year recurring", size=15, color=TEXT_DARK)
add_bullet(tf, "Scales against you as you grow", size=15, color=RED, bold=True)
add_bullet(tf, "No email automation included", size=15, color=TEXT_DARK)
add_bullet(tf, "Data on Atlassian servers", size=15, color=TEXT_DARK)

# RIGHT: FHE (green)
add_rounded(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.2), RGBColor(0xE4, 0xF7, 0xE1))
add_shape(slide, Inches(6.9), Inches(1.6), Inches(5.6), Inches(0.6), GREEN)
add_text(slide, 'FHE Project Board (5 years)', Inches(7.1), Inches(1.65), Inches(5.2), Inches(0.5),
         size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, '$1,275', Inches(7.3), Inches(2.5), Inches(5), Inches(0.8),
         size=56, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

txBox = add_text(slide, '', Inches(7.5), Inches(3.5), Inches(4.5), Inches(2.0), size=15, color=TEXT_DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = "~$255/year hosting + domain"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.color.rgb = TEXT_DARK
add_bullet(tf, "UNLIMITED users - no per-seat cost", size=15, color=GREEN, bold=True)
add_bullet(tf, "Email automation BUILT IN", size=15, color=TEXT_DARK)
add_bullet(tf, "Full data ownership", size=15, color=TEXT_DARK)
add_bullet(tf, "Custom workflows", size=15, color=TEXT_DARK)

# Bottom savings banner
add_rounded(slide, Inches(2.5), Inches(6.1), Inches(8.333), Inches(1.1), DARK_BG)
add_text(slide, 'SAVE $13,725+ OVER 5 YEARS', Inches(2.5), Inches(6.2), Inches(8.333), Inches(0.6),
         size=32, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "That's 91% less than Trello. Math doesn't lie.",
         Inches(2.5), Inches(6.7), Inches(8.333), Inches(0.4),
         size=14, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: ROI BREAKDOWN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.06), PRIMARY)

add_text(slide, "The Math Doesn't Lie", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
         size=38, color=WHITE, bold=True)

# Table data
headers = ['', 'Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5', 'TOTAL']
trello_row = ['Trello (growing to 20)', '$1,050', '$1,800', '$2,400', '$3,000', '$3,000', '$11,250']
fhe_row = ['FHE Project Board', '$500', '$255', '$255', '$255', '$255', '$1,520']
savings_row = ['SAVINGS', '$550', '$1,545', '$2,145', '$2,745', '$2,745', '$9,730']

# Build table
rows, cols = 4, 7
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.4))
table = table_shape.table

# Style header row
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BG

# Trello row (red tint)
for j, val in enumerate(trello_row):
    cell = table.cell(1, j)
    cell.text = val
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = RED if j > 0 else TEXT_DARK
        p.font.bold = j == len(trello_row) - 1
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xFC, 0xEA, 0xE8) if j > 0 else WHITE

# FHE row (green tint)
for j, val in enumerate(fhe_row):
    cell = table.cell(2, j)
    cell.text = val
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = GREEN if j > 0 else TEXT_DARK
        p.font.bold = j == len(fhe_row) - 1
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xE4, 0xF7, 0xE1) if j > 0 else WHITE

# Savings row
for j, val in enumerate(savings_row):
    cell = table.cell(3, j)
    cell.text = val
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = GREEN if j > 0 else RGBColor(0x1B, 0x7A, 0x0A)

# Key callouts
# Break even
box = add_rounded(slide, Inches(0.8), Inches(4.4), Inches(3.6), Inches(1.3), RGBColor(0xE4, 0xF7, 0xE1))
add_text(slide, '\u23F0', Inches(1.0), Inches(4.5), Inches(0.5), Inches(0.5), size=28, color=GREEN)
add_text(slide, 'BREAK-EVEN', Inches(1.5), Inches(4.5), Inches(2.5), Inches(0.3),
         size=14, color=GREEN, bold=True)
add_text(slide, 'Month 3', Inches(1.5), Inches(4.85), Inches(2.5), Inches(0.4),
         size=28, color=TEXT_DARK, bold=True)
add_text(slide, 'Pays for itself in 90 days', Inches(1.5), Inches(5.3), Inches(2.5), Inches(0.3),
         size=12, color=TEXT_SEC)

# Per hire
box = add_rounded(slide, Inches(4.9), Inches(4.4), Inches(3.6), Inches(1.3), RGBColor(0xE0, 0xF4, 0xFD))
add_text(slide, '\U0001f465', Inches(5.1), Inches(4.5), Inches(0.5), Inches(0.5), size=28, color=PRIMARY)
add_text(slide, 'COST PER NEW HIRE', Inches(5.6), Inches(4.5), Inches(2.5), Inches(0.3),
         size=14, color=PRIMARY, bold=True)
add_text(slide, '$0', Inches(5.6), Inches(4.85), Inches(2.5), Inches(0.4),
         size=28, color=TEXT_DARK, bold=True)
add_text(slide, 'vs $12.50/month with Trello', Inches(5.6), Inches(5.3), Inches(2.5), Inches(0.3),
         size=12, color=TEXT_SEC)

# 5 year savings
box = add_rounded(slide, Inches(9.0), Inches(4.4), Inches(3.6), Inches(1.3), RGBColor(0xFF, 0xF4, 0xE5))
add_text(slide, '\U0001f4b0', Inches(9.2), Inches(4.5), Inches(0.5), Inches(0.5), size=28, color=RGBColor(0xD2, 0x90, 0x34))
add_text(slide, '5-YEAR SAVINGS', Inches(9.7), Inches(4.5), Inches(2.5), Inches(0.3),
         size=14, color=RGBColor(0xD2, 0x90, 0x34), bold=True)
add_text(slide, '$9,730+', Inches(9.7), Inches(4.85), Inches(2.5), Inches(0.4),
         size=28, color=TEXT_DARK, bold=True)
add_text(slide, 'Conservative estimate', Inches(9.7), Inches(5.3), Inches(2.5), Inches(0.3),
         size=12, color=TEXT_SEC)

# Bottom note
add_text(slide, '"Every new hire is FREE, not $12.50/month. Growth should make us money, not cost us more in software."',
         Inches(1), Inches(6.3), Inches(11.333), Inches(0.5),
         size=15, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: BONUS FEATURES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_TEAL)

add_text(slide, "But Wait, There's More!", Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
         size=38, color=WHITE, bold=True)
add_text(slide, "(We promise this isn\'t an infomercial)",
         Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
         size=16, color=TEXT_SEC)

bonuses = [
    ("\u2709 Email Automation", "Incoming emails automatically become task cards.\nFilter by sender, subject, recipient.\nNo more manual data entry."),
    ("\u2699 Custom Workflows", "Tailored to how FHE actually works.\nNot forced into Trello's templates.\nOur processes, our rules."),
    ("\U0001f512 Full Data Ownership", "Data stays on OUR servers.\nNo third-party access. Full compliance.\nBackup on our schedule."),
    ("\U0001f4f1 Works Everywhere", "Mobile responsive PWA.\nDesktop, tablet, phone.\nOffline support built in."),
]

for i, (title, desc) in enumerate(bonuses):
    col = i % 2
    row = i // 2
    x = 0.6 + col * 6.4
    y = 1.7 + row * 2.6

    card = add_rounded(slide, Inches(x), Inches(y), Inches(5.9), Inches(2.2), RGBColor(0x14, 0x28, 0x3C))

    add_text(slide, title, Inches(x + 0.4), Inches(y + 0.2), Inches(5), Inches(0.4),
             size=22, color=ACCENT_TEAL, bold=True)

    add_text(slide, desc, Inches(x + 0.4), Inches(y + 0.75), Inches(5), Inches(1.3),
             size=15, color=TEXT_SEC)

# Bottom zinger
add_text(slide, '"No more \'Sorry, that\'s a Premium feature.\' Everything is a Premium feature. It\'s ours."',
         Inches(1), Inches(6.7), Inches(11.333), Inches(0.4),
         size=16, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: THE ASK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.06), PRIMARY)

add_text(slide, 'The Ask', Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
         size=38, color=WHITE, bold=True)

# Three columns for the ask
asks = [
    ("\u23F1", "Development Time", "40-60 hours", "to finish & polish", PRIMARY),
    ("\U0001f4bb", "Hosting Budget", "$20/month", "cloud hosting", GREEN),
    ("\U0001f4c5", "Timeline", "2-4 weeks", "to full deployment", ACCENT_TEAL),
]

for i, (icon, title, big, sub, color) in enumerate(asks):
    x = 1.0 + i * 4.1
    card = add_rounded(slide, Inches(x), Inches(1.8), Inches(3.6), Inches(3.0), LIGHT_BG)
    add_shape(slide, Inches(x), Inches(1.8), Inches(3.6), Inches(0.08), color)

    add_text(slide, icon, Inches(x), Inches(2.1), Inches(3.6), Inches(0.5),
             size=36, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(x), Inches(2.65), Inches(3.6), Inches(0.3),
             size=14, color=TEXT_SEC, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, big, Inches(x), Inches(3.1), Inches(3.6), Inches(0.5),
             size=32, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, sub, Inches(x), Inches(3.65), Inches(3.6), Inches(0.3),
             size=14, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# Perspective box
box = add_rounded(slide, Inches(1.5), Inches(5.2), Inches(10.333), Inches(1.6), RGBColor(0xE4, 0xF7, 0xE1))

add_text(slide, 'To put it in perspective:', Inches(2.0), Inches(5.35), Inches(9), Inches(0.3),
         size=16, color=TEXT_DARK, bold=True)

txBox = add_text(slide, '', Inches(2.0), Inches(5.7), Inches(9), Inches(1.0), size=16, color=TEXT_DARK)
tf = txBox.text_frame
tf.paragraphs[0].text = "Total investment is less than ONE MONTH of Trello at scale ($250)"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.color.rgb = TEXT_DARK
add_bullet(tf, "We already have a working prototype (you just saw it)", size=16, color=GREEN, bold=True)
add_bullet(tf, "ROI is positive by Month 3", size=16, color=TEXT_DARK)

# ============================================================
# SLIDE 10: CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

# Accent bars
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), GREEN)

# Decorative elements
c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.8), Inches(5.5), Inches(2.5), Inches(2.5))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x01, 0x45, 0x70); c1.line.fill.background()
c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(-0.5), Inches(2.5), Inches(2.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x01, 0x45, 0x70); c2.line.fill.background()

add_text(slide, 'Own Your Tools.', Inches(1), Inches(1.5), Inches(11.333), Inches(0.9),
         size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 'Own Your Future.', Inches(1), Inches(2.4), Inches(11.333), Inches(0.9),
         size=52, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Divider
add_shape(slide, Inches(5.5), Inches(3.6), Inches(2.333), Inches(0.04), PRIMARY)

add_text(slide, 'FHE Project Board', Inches(1), Inches(3.9), Inches(11.333), Inches(0.6),
         size=28, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 'Built by Engineers, for Engineers', Inches(1), Inches(4.5), Inches(11.333), Inches(0.4),
         size=18, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# Summary boxes
summary = [
    ("91% Cost Savings", GREEN),
    ("Unlimited Users", PRIMARY),
    ("Email Automation", ACCENT_TEAL),
    ("Full Data Control", RGBColor(0x89, 0x60, 0x9E)),
]

total_width = len(summary) * 2.6 + (len(summary)-1) * 0.3
start_x = (13.333 - total_width) / 2
for i, (text, color) in enumerate(summary):
    x = start_x + i * 2.9
    box = add_rounded(slide, Inches(x), Inches(5.2), Inches(2.6), Inches(0.7), color)
    add_text(slide, text, Inches(x), Inches(5.3), Inches(2.6), Inches(0.5),
             size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Closing zinger
add_text(slide, 'P.S. Atlassian won\'t miss us. They have plenty of subscribers.',
         Inches(1), Inches(6.5), Inches(11.333), Inches(0.4),
         size=16, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# ============ SAVE ============
output_path = '/home/user/fhe-TRELLO-/FHE_Project_Board_Pitch.pptx'
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"File size: {os.path.getsize(output_path) / 1024:.1f} KB")
print(f"Total slides: {len(prs.slides)}")
